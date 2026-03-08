import asyncio
import logging
import os
import re
import uuid
from datetime import date, datetime

import httpx
from sqlalchemy import select, or_, func, text
from sqlalchemy.ext.asyncio import AsyncSession

from app.config import settings
from app.database import async_session
from app.models.document import OpportunityDocument, DocumentChunk
from app.models.opportunity import Opportunity
from app.services.grants_client import GrantsGovClient
from app.services.settings_service import settings_service
from app.services.cache_service import cache_service

logger = logging.getLogger(__name__)

REDIS_DOC_SYNC_KEY = "pf:doc_sync_stats"
REDIS_DOC_PROCESSING_FLAG = "pf:doc_processing"
REDIS_DOC_COMPLETED_KEY = "pf:doc_completed"


class DocumentService:
    def __init__(self):
        self.is_processing = False
        self._cancel_requested = False
        self.processing_stats: dict = {}
        self._grants_client = GrantsGovClient()

    async def extract_attachment_metadata(
        self, session: AsyncSession, opp, detail: dict
    ) -> int:
        """Parse attachment metadata from fetchOpportunity response and upsert rows.

        Returns the number of new documents created.
        """
        folders = detail.get("synopsisAttachmentFolders") or []
        created = 0

        for folder in folders:
            folder_name = folder.get("folderName", "")
            attachments = folder.get("synopsisAttachments") or []

            for att in attachments:
                att_id = str(att.get("id", ""))
                if not att_id:
                    continue

                # Check if already exists
                stmt = select(OpportunityDocument).where(
                    OpportunityDocument.attachment_id == att_id
                )
                result = await session.execute(stmt)
                existing = result.scalar_one_or_none()

                if existing:
                    # Update metadata if changed
                    existing.file_name = att.get("fileName", existing.file_name)
                    existing.mime_type = att.get("mimeType", existing.mime_type)
                    existing.file_size = att.get("fileLobSize", existing.file_size)
                    existing.file_description = att.get("fileDescription", existing.file_description)
                    existing.folder_name = folder_name
                    continue

                doc = OpportunityDocument(
                    opportunity_id=opp.id,
                    attachment_id=att_id,
                    file_name=att.get("fileName", "unknown"),
                    mime_type=att.get("mimeType"),
                    file_size=att.get("fileLobSize"),
                    file_description=att.get("fileDescription"),
                    folder_name=folder_name,
                    download_status="pending",
                    ocr_status="pending",
                    embed_status="pending",
                )
                session.add(doc)
                created += 1

        if created > 0:
            await session.flush()

        return created

    # Domains to skip when extracting linked documents
    _SKIP_DOMAINS = {
        "sam.gov", "www.sam.gov",
        "grants.gov", "www.grants.gov", "apply07.grants.gov",
        "teams.microsoft.com", "dod.teams.microsoft.us",
        "youtube.com", "www.youtube.com",
        "twitter.com", "x.com",
        "facebook.com", "www.facebook.com",
        "linkedin.com", "www.linkedin.com",
        "whitehouse.gov", "www.whitehouse.gov",
        "urldefense.proofpoint.com",
        "usaspending.gov", "www.usaspending.gov",
        "sba.gov", "www.sba.gov",
    }

    # URL path patterns to skip entirely (noise pages)
    _SKIP_URL_PATTERNS = re.compile(
        r'/(contact|about|staff|careers|login|signup|faq|blog|news|press|privacy|terms|accessibility|sitemap|help)(/|$)',
        re.IGNORECASE,
    )

    @staticmethod
    def _score_url(url: str) -> int:
        """Score a URL by how likely it is to be a useful document. Negative = skip."""
        lower = url.lower()
        path = lower.split("?")[0]  # ignore query string for scoring

        # Filter out noise pages
        if DocumentService._SKIP_URL_PATTERNS.search(path):
            return -1

        score = 0

        # File extension bonuses
        if path.endswith(".pdf"):
            score += 50
        elif path.endswith(".docx") or path.endswith(".doc"):
            score += 40

        # Solicitation keywords in path
        sol_keywords = ("nofo", "rfp", "rfa", "foa", "solicitation", "baa", "announcement")
        for kw in sol_keywords:
            if kw in path:
                score += 30
                break

        # Document-related keywords
        doc_keywords = ("download", "attachment", "file", "apply", "document")
        for kw in doc_keywords:
            if kw in path:
                score += 15
                break

        # Penalize browse/search/listing pages
        browse_keywords = ("browse", "search", "list", "archive", "index", "category")
        for kw in browse_keywords:
            if kw in path:
                score -= 10
                break

        return score

    # Content types we can process
    _FETCHABLE_CONTENT_TYPES = {
        "application/pdf", "text/html",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # docx
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # xlsx
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # pptx
        "application/msword",  # doc
        "text/plain", "text/csv",
    }

    _URL_PATTERN = re.compile(r'https?://[^\s<>"\')\]&;]+(?:&amp;[^\s<>"\')\]&;]+)*')

    def _extract_urls(self, text_content: str) -> list[str]:
        """Extract and deduplicate URLs from text, filtering noise domains."""
        import html
        from urllib.parse import urlparse

        raw_urls = self._URL_PATTERN.findall(text_content)
        seen = set()
        result = []

        for url in raw_urls:
            # Decode HTML entities
            url = html.unescape(url)
            # Strip trailing punctuation that's not part of URL
            url = url.rstrip(".,;:!?)")

            try:
                parsed = urlparse(url)
                domain = parsed.netloc.lower().rstrip(".")
            except Exception:
                continue

            if not domain:
                continue
            if domain in self._SKIP_DOMAINS:
                continue
            if url in seen:
                continue

            seen.add(url)
            result.append(url)

        return result

    async def extract_linked_documents(
        self, session: AsyncSession, opp, client: httpx.AsyncClient | None = None
    ) -> int:
        """Extract URLs from opportunity description, fetch content, create document rows.

        Returns the number of new linked documents created.
        """
        import hashlib
        from urllib.parse import urlparse

        description = opp.synopsis_description or ""
        urls = self._extract_urls(description)
        if not urls:
            return 0

        # Score, filter, sort, and limit URLs
        scored = [(self._score_url(u), u) for u in urls]
        scored = [(s, u) for s, u in scored if s >= 0]
        scored.sort(key=lambda x: x[0], reverse=True)
        urls = [u for _, u in scored[:5]]
        if not urls:
            return 0

        created = 0
        close_client = False
        if client is None:
            client = httpx.AsyncClient(
                follow_redirects=True,
                timeout=30.0,
                headers={"User-Agent": "ProposalForge/1.0 (Federal Grant Discovery)"},
            )
            close_client = True

        try:
            for url in urls:
                url_hash = hashlib.md5(url.encode()).hexdigest()[:12]
                att_id = f"linked_{url_hash}"

                # Skip if already exists
                stmt = select(OpportunityDocument).where(
                    OpportunityDocument.attachment_id == att_id
                )
                result = await session.execute(stmt)
                if result.scalar_one_or_none():
                    continue

                try:
                    resp = await client.get(url)
                    resp.raise_for_status()
                except Exception as e:
                    logger.debug(f"Failed to fetch {url}: {e}")
                    continue

                content_type = (resp.headers.get("content-type") or "").split(";")[0].strip().lower()
                content_length = len(resp.content)

                # Skip if too large (>10MB) or empty
                if content_length > 10 * 1024 * 1024 or content_length == 0:
                    continue

                # Skip unsupported content types
                if not any(content_type.startswith(ct) for ct in self._FETCHABLE_CONTENT_TYPES):
                    continue

                # Determine file extension from content type
                ext_map = {
                    "application/pdf": ".pdf",
                    "text/html": ".html",
                    "text/plain": ".txt",
                    "text/csv": ".csv",
                    "application/msword": ".doc",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
                }
                ext = ext_map.get(content_type, "")

                # Try to get filename from URL path
                try:
                    path = urlparse(str(resp.url)).path
                    url_filename = path.split("/")[-1] if path else ""
                    if not url_filename or len(url_filename) > 200:
                        url_filename = ""
                except Exception:
                    url_filename = ""

                if url_filename and "." in url_filename:
                    file_name = url_filename
                else:
                    file_name = f"linked_{url_hash}{ext}"

                safe_name = re.sub(r'[^\w\-.]', '_', file_name)
                dest_path = os.path.join(
                    settings.DOCUMENT_STORAGE_PATH,
                    str(opp.id),
                    f"{att_id}_{safe_name}",
                )
                os.makedirs(os.path.dirname(dest_path), exist_ok=True)

                with open(dest_path, "wb") as f:
                    f.write(resp.content)

                doc = OpportunityDocument(
                    opportunity_id=opp.id,
                    attachment_id=att_id,
                    file_name=file_name,
                    mime_type=content_type,
                    file_size=content_length,
                    file_description=f"Linked from description: {url[:500]}",
                    folder_name="Linked from description",
                    source="linked",
                    local_path=dest_path,
                    download_status="downloaded",
                    ocr_status="pending",
                    classify_status="pending",
                    embed_status="pending",
                )
                session.add(doc)
                created += 1

        finally:
            if close_client:
                await client.aclose()

        if created > 0:
            await session.flush()

        return created

    async def batch_extract_linked_documents(self):
        """Batch extract linked documents for all open opportunities without Grants.gov docs."""
        if self.is_processing:
            logger.warning("Document processing already in progress, skipping link extraction")
            return 0

        self.is_processing = True
        self._cancel_requested = False
        self.processing_stats = {
            "started": datetime.utcnow().isoformat(),
            "total": 0,
            "scanned": 0,
            "downloaded": 0,
            "errors": 0,
            "phase": "extracting links",
        }
        await self._publish_stats()

        try:
            from app.models.opportunity import Opportunity

            async with async_session() as session:
                today = date.today()
                # Find open opportunities with no documents at all
                stmt = (
                    select(Opportunity)
                    .outerjoin(OpportunityDocument, OpportunityDocument.opportunity_id == Opportunity.id)
                    .where(
                        Opportunity.status != "archived",
                        or_(Opportunity.close_date >= today, Opportunity.close_date.is_(None)),
                        Opportunity.synopsis_description.isnot(None),
                        OpportunityDocument.id.is_(None),
                    )
                )
                result = await session.execute(stmt)
                opps = result.scalars().unique().all()

                self.processing_stats["total"] = len(opps)
                await self._publish_stats()
                logger.info(f"Batch link extraction: {len(opps)} opportunities to scan")
                total_created = 0

                async with httpx.AsyncClient(
                    follow_redirects=True,
                    timeout=30.0,
                    headers={"User-Agent": "ProposalForge/1.0 (Federal Grant Discovery)"},
                ) as client:
                    for i, opp in enumerate(opps):
                        if self._cancel_requested:
                            logger.info("Link extraction cancelled")
                            self.processing_stats["phase"] = "cancelled"
                            break

                        try:
                            created = await self.extract_linked_documents(session, opp, client=client)
                            if created > 0:
                                await session.commit()
                                total_created += created
                                self.processing_stats["downloaded"] = total_created
                                logger.info(f"  [{i+1}/{len(opps)}] {opp.opportunity_id}: {created} linked docs")
                        except Exception as e:
                            logger.warning(f"  [{i+1}/{len(opps)}] {opp.opportunity_id}: error - {e}")
                            self.processing_stats["errors"] = self.processing_stats.get("errors", 0) + 1
                            await session.rollback()

                        self.processing_stats["scanned"] = i + 1
                        if i % 5 == 4:
                            await self._publish_stats()

                        # Rate limit
                        if i % 10 == 9:
                            await asyncio.sleep(1)

                if self.processing_stats["phase"] != "cancelled":
                    self.processing_stats["phase"] = "completed"
                self.processing_stats["completed"] = datetime.utcnow().isoformat()
                logger.info(f"Batch link extraction complete: {total_created} total docs created")

        except Exception as e:
            logger.error(f"Batch link extraction failed: {e}", exc_info=True)
            self.processing_stats["error"] = str(e)[:500]

        # Create synopsis docs for remaining doc-less opportunities
        if not self._cancel_requested:
            try:
                self.is_processing = True  # ensure flag is set
                await self._create_synopsis_documents_inline()
            except Exception as e:
                logger.error(f"Synopsis doc creation failed: {e}", exc_info=True)

        self.is_processing = False
        await self._publish_stats()

        # Chain into document processing, skip link extraction since we just did it
        if not self._cancel_requested:
            logger.info("Chaining into document processing pipeline...")
            await self.process_pending_documents(skip_link_extraction=True)

    async def _inline_link_extraction(self):
        """Scan open opportunities without docs and extract linked documents.

        Called inline during process_pending_documents (self.is_processing already True).
        Only scans opportunities whose synopsis actually contains URLs (LIKE '%http%').
        """
        from app.models.opportunity import Opportunity

        self.processing_stats["phase"] = "extracting links"
        await self._publish_stats()

        async with async_session() as session:
            today = date.today()
            stmt = (
                select(Opportunity)
                .outerjoin(OpportunityDocument, OpportunityDocument.opportunity_id == Opportunity.id)
                .where(
                    Opportunity.status != "archived",
                    or_(Opportunity.close_date >= today, Opportunity.close_date.is_(None)),
                    Opportunity.synopsis_description.isnot(None),
                    Opportunity.synopsis_description.like("%http%"),
                    OpportunityDocument.id.is_(None),
                )
            )
            result = await session.execute(stmt)
            opps = result.scalars().unique().all()

            if not opps:
                logger.info("No opportunities without documents to scan for links")
                return

            self.processing_stats["total"] = len(opps)
            self.processing_stats["scanned"] = 0
            self.processing_stats["downloaded"] = 0
            await self._publish_stats()
            logger.info(f"Link extraction: scanning {len(opps)} opportunities (with URLs)")

            total_created = 0
            fetch_semaphore = asyncio.Semaphore(10)

            async with httpx.AsyncClient(
                follow_redirects=True,
                timeout=30.0,
                headers={"User-Agent": "ProposalForge/1.0 (Federal Grant Discovery)"},
            ) as client:
                batch_size = 100
                for batch_start in range(0, len(opps), batch_size):
                    if self._cancel_requested:
                        break
                    batch = opps[batch_start:batch_start + batch_size]

                    async def _process_opp(idx, opp):
                        nonlocal total_created
                        if self._cancel_requested:
                            return
                        async with fetch_semaphore:
                            if self._cancel_requested:
                                return
                            try:
                                created = await self.extract_linked_documents(session, opp, client=client)
                                if created > 0:
                                    await session.commit()
                                    total_created += created
                                    self.processing_stats["downloaded"] = total_created
                            except Exception as e:
                                logger.warning(f"Link extraction for {opp.opportunity_id}: {e}")
                                self.processing_stats["errors"] = self.processing_stats.get("errors", 0) + 1
                                await session.rollback()
                            self.processing_stats["scanned"] = batch_start + idx + 1

                    await asyncio.gather(*[_process_opp(i, opp) for i, opp in enumerate(batch)])
                    await self._publish_stats()

            logger.info(f"Link extraction complete: {total_created} docs from {len(opps)} opportunities")

        # After link extraction, create synopsis docs for remaining doc-less opportunities
        if not self._cancel_requested:
            await self._create_synopsis_documents_inline()

    async def _create_synopsis_documents_inline(self):
        """Create synthetic documents from synopsis descriptions for opportunities still without docs.

        Called inline during _inline_link_extraction (self.is_processing already True).
        """
        self.processing_stats["phase"] = "creating synopsis docs"
        await self._publish_stats()

        async with async_session() as session:
            today = date.today()
            # Find open opportunities with no documents — include those with short descriptions too
            stmt = (
                select(Opportunity)
                .outerjoin(OpportunityDocument, OpportunityDocument.opportunity_id == Opportunity.id)
                .where(
                    Opportunity.status != "archived",
                    or_(Opportunity.close_date >= today, Opportunity.close_date.is_(None)),
                    OpportunityDocument.id.is_(None),
                )
            )
            result = await session.execute(stmt)
            opps = result.scalars().unique().all()

            if not opps:
                logger.info("No opportunities need synopsis documents")
                return

            logger.info(f"Creating synopsis documents for {len(opps)} opportunities")
            created = 0

            for opp in opps:
                if self._cancel_requested:
                    break
                try:
                    # Build a stable attachment_id
                    att_id = f"synopsis_{opp.id}"

                    # Check if already exists
                    existing = (await session.execute(
                        select(OpportunityDocument).where(OpportunityDocument.attachment_id == att_id)
                    )).scalar_one_or_none()
                    if existing:
                        continue

                    synopsis = (opp.synopsis_description or "").strip()

                    # If synopsis is too short or empty, build metadata-only synopsis
                    if len(synopsis) <= 50:
                        parts = []
                        if opp.title:
                            parts.append(f"Title: {opp.title}")
                        if opp.opportunity_number:
                            parts.append(f"Opportunity Number: {opp.opportunity_number}")
                        if opp.agency_code:
                            parts.append(f"Agency: {opp.agency_code}")
                        if getattr(opp, "award_floor", None) or getattr(opp, "award_ceiling", None):
                            floor = getattr(opp, "award_floor", None)
                            ceiling = getattr(opp, "award_ceiling", None)
                            if floor and ceiling:
                                parts.append(f"Funding Range: ${floor:,.0f} - ${ceiling:,.0f}")
                            elif ceiling:
                                parts.append(f"Funding Up To: ${ceiling:,.0f}")
                        if opp.close_date:
                            parts.append(f"Close Date: {opp.close_date}")
                        if not parts:
                            continue  # Nothing useful to write
                        synopsis = "\n".join(parts)

                    # Write synopsis text to disk
                    opp_dir = os.path.join(settings.DOCUMENT_STORAGE_PATH, str(opp.id))
                    os.makedirs(opp_dir, exist_ok=True)
                    dest_path = os.path.join(opp_dir, "Synopsis_Description.txt")
                    text_path = dest_path + ".txt"  # Companion extracted text file

                    with open(dest_path, "w", encoding="utf-8") as f:
                        f.write(synopsis)
                    # Write identical extracted text companion
                    with open(text_path, "w", encoding="utf-8") as f:
                        f.write(synopsis)

                    doc = OpportunityDocument(
                        opportunity_id=opp.id,
                        attachment_id=att_id,
                        file_name="Synopsis Description.txt",
                        mime_type="text/plain",
                        file_size=len(synopsis.encode("utf-8")),
                        file_description="Opportunity synopsis description (auto-generated)",
                        folder_name="Synopsis",
                        source="synopsis",
                        local_path=dest_path,
                        download_status="downloaded",
                        ocr_status="completed",
                        classify_status="pending",
                        embed_status="pending",
                        extracted_text_length=len(synopsis),
                    )
                    session.add(doc)
                    created += 1

                    if created % 100 == 0:
                        await session.flush()

                except Exception as e:
                    logger.warning(f"Synopsis doc for opportunity {opp.id}: {e}")
                    await session.rollback()

            if created > 0:
                await session.commit()
            logger.info(f"Created {created} synopsis documents")

    async def create_synopsis_documents(self):
        """Standalone batch: create synopsis documents for all doc-less opportunities."""
        if self.is_processing:
            logger.warning("Document processing already in progress")
            return 0

        self.is_processing = True
        self._cancel_requested = False
        self.processing_stats = {
            "started": datetime.utcnow().isoformat(),
            "phase": "creating synopsis docs",
            "total": 0, "created": 0, "errors": 0,
        }
        await self._publish_stats()

        try:
            await self._create_synopsis_documents_inline()
        finally:
            self.is_processing = False
            self.processing_stats["phase"] = "completed"
            self.processing_stats["completed"] = datetime.utcnow().isoformat()
            await self._publish_stats()

    # --- Agency domain mapping for web search ---
    _AGENCY_DOMAINS = {
        "NASA": ["nspires.nasaprs.com", "nasa.gov"],
        "NSF": ["nsf.gov"],
        "DOE": ["energy.gov", "science.energy.gov"],
        "DOD": ["grants.darpa.mil", "arl.army.mil", "defense.gov"],
        "HHS": ["grants.nih.gov", "hrsa.gov", "samhsa.gov", "hhs.gov"],
        "NIH": ["grants.nih.gov", "nih.gov"],
        "EPA": ["epa.gov"],
        "USDA": ["nifa.usda.gov", "usda.gov"],
        "ED": ["ed.gov"],
        "DOJ": ["ojp.gov", "justice.gov"],
        "DOT": ["transportation.gov"],
        "DOC": ["commerce.gov", "nist.gov", "noaa.gov"],
        "DHS": ["dhs.gov", "fema.gov"],
        "DOI": ["doi.gov", "fws.gov", "nps.gov"],
        "DARPA": ["darpa.mil"],
    }

    async def search_for_solicitations(self):
        """Web search for solicitation PDFs for opportunities lacking solicitation docs."""
        if self.is_processing:
            logger.warning("Document processing already in progress")
            return 0

        if not settings.BRAVE_API_KEY:
            logger.warning("BRAVE_API_KEY not configured, cannot run web search")
            return 0

        self.is_processing = True
        self._cancel_requested = False
        total_found = 0
        self.processing_stats = {
            "started": datetime.utcnow().isoformat(),
            "phase": "web search",
            "total": 0, "scanned": 0, "downloaded": 0, "errors": 0,
        }
        await self._publish_stats()

        try:
            async with async_session() as session:
                today = date.today()
                # Find open opportunities that have no solicitation document
                from sqlalchemy.orm import aliased
                SolDoc = aliased(OpportunityDocument)
                # Subquery: opportunities that already have a solicitation doc
                has_sol = (
                    select(SolDoc.opportunity_id)
                    .where(SolDoc.doc_category == "solicitation")
                    .correlate(Opportunity)
                )
                stmt = (
                    select(Opportunity)
                    .where(
                        Opportunity.status.notin_(["archived", "closed"]),
                        or_(Opportunity.close_date >= today, Opportunity.close_date.is_(None)),
                        ~Opportunity.id.in_(has_sol),
                    )
                )
                result = await session.execute(stmt)
                opps = result.scalars().all()

                self.processing_stats["total"] = len(opps)
                await self._publish_stats()
                logger.info(f"Web search: {len(opps)} opportunities to search")

                total_found = 0
                async with httpx.AsyncClient(
                    follow_redirects=True,
                    timeout=30.0,
                    headers={"User-Agent": "ProposalForge/1.0 (Federal Grant Discovery)"},
                ) as client:
                    for i, opp in enumerate(opps):
                        if self._cancel_requested:
                            break

                        try:
                            found = await self._web_search_for_opportunity(opp, session, client)
                            if found > 0:
                                await session.commit()
                                total_found += found
                                self.processing_stats["downloaded"] = total_found
                        except Exception as e:
                            logger.warning(f"Web search for {opp.opportunity_id}: {e}")
                            self.processing_stats["errors"] = self.processing_stats.get("errors", 0) + 1
                            await session.rollback()

                        self.processing_stats["scanned"] = i + 1
                        if i % 5 == 4:
                            await self._publish_stats()
                        # Rate limit: respect Brave API limits
                        await asyncio.sleep(1)

                self.processing_stats["phase"] = "completed"
                self.processing_stats["completed"] = datetime.utcnow().isoformat()
                logger.info(f"Web search complete: {total_found} solicitations found")

        except Exception as e:
            logger.error(f"Web search failed: {e}", exc_info=True)
            self.processing_stats["error"] = str(e)[:500]
        finally:
            self.is_processing = False
            await self._publish_stats()

        # Chain into document processing to handle newly downloaded docs
        if not self._cancel_requested and total_found > 0:
            logger.info("Chaining into document processing pipeline...")
            await self.process_pending_documents(skip_link_extraction=True)

        return total_found

    async def _web_search_for_opportunity(
        self, opp, session: AsyncSession, client: httpx.AsyncClient
    ) -> int:
        """Search the web for solicitation PDFs for a single opportunity."""
        title = (opp.title or "").strip()
        opp_number = (opp.opportunity_number or "").strip()
        agency_code = (opp.agency_code or "").upper()

        if not title and not opp_number:
            return 0

        # Build search query — keep it short to avoid Brave 422 errors
        query_parts = []
        if opp_number:
            query_parts.append(f'"{opp_number}"')
        if title:
            # Strip parenthetical suffixes like (R21), (R01) and truncate
            clean_title = re.sub(r'\s*\([^)]*\)\s*$', '', title).strip()
            if len(clean_title) > 80:
                clean_title = clean_title[:80].rsplit(' ', 1)[0]
            if clean_title:
                query_parts.append(f'"{clean_title}"')
        query_parts.append("filetype:pdf")

        # Add agency domain hint if available
        for prefix, domains in self._AGENCY_DOMAINS.items():
            if agency_code.startswith(prefix):
                if domains:
                    query_parts.append(f"site:{domains[0]}")
                break

        query = " ".join(query_parts)
        # Cap total query length to avoid Brave API rejections
        if len(query) > 300:
            query = query[:300].rsplit(' ', 1)[0]

        # Call Brave Search API
        try:
            resp = await client.get(
                "https://api.search.brave.com/res/v1/web/search",
                params={"q": query, "count": 5},
                headers={
                    "Accept": "application/json",
                    "Accept-Encoding": "gzip",
                    "X-Subscription-Token": settings.BRAVE_API_KEY,
                },
            )
            resp.raise_for_status()
            data = resp.json()
        except Exception as e:
            logger.warning(f"Brave search failed for {opp.opportunity_id}: {e}")
            return 0

        # Filter for PDF results
        results = data.get("web", {}).get("results", [])
        pdf_urls = []
        for r in results:
            url = r.get("url", "")
            if url.lower().endswith(".pdf"):
                pdf_urls.append(url)
            elif "pdf" in r.get("meta_url", {}).get("path", "").lower():
                pdf_urls.append(url)

        if not pdf_urls:
            return 0

        # Download first matching PDF
        created = 0
        for url in pdf_urls[:2]:  # Try at most 2 URLs
            try:
                # Generate unique attachment ID from URL
                url_hash = str(uuid.uuid5(uuid.NAMESPACE_URL, url))[:16]
                att_id = f"web_{opp.id}_{url_hash}"

                # Check if already exists
                existing = (await session.execute(
                    select(OpportunityDocument).where(OpportunityDocument.attachment_id == att_id)
                )).scalar_one_or_none()
                if existing:
                    continue

                # Download the PDF
                pdf_resp = await client.get(url, timeout=60.0)
                pdf_resp.raise_for_status()

                content_type = pdf_resp.headers.get("content-type", "")
                if "pdf" not in content_type and not url.lower().endswith(".pdf"):
                    continue

                content = pdf_resp.content
                if len(content) < 1000:  # Skip tiny files
                    continue
                if len(content) > 50_000_000:  # Skip >50MB files
                    continue

                # Save to disk
                opp_dir = os.path.join(settings.DOCUMENT_STORAGE_PATH, str(opp.id))
                os.makedirs(opp_dir, exist_ok=True)

                # Extract filename from URL
                from urllib.parse import urlparse, unquote
                url_path = urlparse(url).path
                file_name = unquote(url_path.split("/")[-1]) or "solicitation.pdf"

                dest_path = os.path.join(opp_dir, f"web_{file_name}")

                with open(dest_path, "wb") as f:
                    f.write(content)

                doc = OpportunityDocument(
                    opportunity_id=opp.id,
                    attachment_id=att_id,
                    file_name=file_name,
                    mime_type="application/pdf",
                    file_size=len(content),
                    file_description=f"Found via web search: {url[:200]}",
                    folder_name="Web Search",
                    source="web_search",
                    local_path=dest_path,
                    download_status="downloaded",
                    ocr_status="pending",
                    classify_status="pending",
                    embed_status="pending",
                )
                session.add(doc)
                created += 1
                logger.info(f"  Web search found PDF for {opp.opportunity_id}: {file_name}")
                break  # One solicitation is enough

            except Exception as e:
                logger.debug(f"Failed to download {url}: {e}")
                continue

        return created

    async def process_pending_documents(self, skip_link_extraction: bool = False):
        """Batch orchestrator: extract links, download, OCR, chunk, embed all pending documents."""
        if self.is_processing:
            logger.warning("Document processing already in progress")
            return

        self.is_processing = True
        self._cancel_requested = False
        self.processing_stats = {
            "started": datetime.utcnow().isoformat(),
            "total": 0,
            "downloaded": 0,
            "ocr_completed": 0,
            "classified": 0,
            "embedded": 0,
            "errors": 0,
            "phase": "starting",
        }
        await self._publish_stats()

        try:
            # Phase 0: Extract linked documents from descriptions
            if not skip_link_extraction and not self._cancel_requested:
                await self._inline_link_extraction()

            async with async_session() as session:
                # Get settings once
                ocr_settings = await settings_service.get_ocr_settings(session)
                embed_settings = await settings_service.get_embedding_settings(session)
                llm_settings = await settings_service.get_llm_settings(session)

                # Reset all failed docs back to pending and clear errors
                await session.execute(
                    text("UPDATE opportunity_documents SET download_status = 'pending', error_message = NULL WHERE download_status = 'failed'")
                )
                await session.execute(
                    text("UPDATE opportunity_documents SET ocr_status = 'pending', error_message = NULL WHERE ocr_status = 'failed'")
                )
                await session.execute(
                    text("UPDATE opportunity_documents SET embed_status = 'pending', error_message = NULL WHERE embed_status = 'failed'")
                )
                await session.execute(
                    text("UPDATE opportunity_documents SET classify_status = 'pending', error_message = NULL WHERE classify_status = 'failed'")
                )
                await session.commit()
                logger.info("Reset failed documents to pending")

                # Delete documents for closed/archived opportunities
                today = date.today()
                closed_docs = await session.execute(
                    select(OpportunityDocument.id)
                    .join(Opportunity, OpportunityDocument.opportunity_id == Opportunity.id)
                    .where(or_(
                        Opportunity.status == "archived",
                        Opportunity.close_date < today,
                    ))
                )
                closed_ids = [row[0] for row in closed_docs.all()]
                if closed_ids:
                    for i in range(0, len(closed_ids), 500):
                        batch = closed_ids[i:i + 500]
                        await session.execute(
                            DocumentChunk.__table__.delete().where(
                                DocumentChunk.document_id.in_(batch)
                            )
                        )
                        await session.execute(
                            OpportunityDocument.__table__.delete().where(
                                OpportunityDocument.id.in_(batch)
                            )
                        )
                        await session.commit()
                    logger.info(f"Deleted {len(closed_ids)} documents from closed/archived opportunities")

                # Query pending docs for open opportunities only
                stmt = (
                    select(OpportunityDocument)
                    .join(Opportunity, OpportunityDocument.opportunity_id == Opportunity.id)
                    .where(
                        Opportunity.status != "archived",
                        or_(Opportunity.close_date >= today, Opportunity.close_date.is_(None)),
                        or_(
                            OpportunityDocument.download_status == "pending",
                            OpportunityDocument.ocr_status == "pending",
                            OpportunityDocument.classify_status == "pending",
                            OpportunityDocument.embed_status == "pending",
                        ),
                    )
                )
                result = await session.execute(stmt)
                docs = list(result.scalars().all())

            num_workers = ocr_settings.get("doc_workers", 4)
            self.processing_stats["total"] = len(docs)
            self.processing_stats["phase"] = "processing"
            await self._publish_stats()

            logger.info(f"Processing {len(docs)} pending documents with {num_workers} workers")

            # Create shared clients to avoid exhausting file descriptors
            import chromadb
            from openai import AsyncOpenAI

            base_url = embed_settings.get("base_url", "")
            model = embed_settings.get("model", "")
            api_key = embed_settings.get("api_key", "")
            shared_embed_client = AsyncOpenAI(base_url=base_url, api_key=api_key or "not-needed") if base_url and model else None

            llm_base = llm_settings.get("base_url", "")
            llm_model = llm_settings.get("model", "")
            llm_key = llm_settings.get("api_key", "")
            shared_llm_client = AsyncOpenAI(base_url=llm_base, api_key=llm_key or "not-needed") if llm_base and llm_model else None
            shared_chroma_client = chromadb.HttpClient(
                host=settings.CHROMADB_HOST,
                port=settings.CHROMADB_PORT,
            )
            shared_chroma_collection = shared_chroma_client.get_or_create_collection(
                name="opportunity_documents",
                metadata={"hnsw:space": "cosine"},
            )
            shared_ocr_client = httpx.AsyncClient(timeout=300.0)

            semaphore = asyncio.Semaphore(num_workers)

            async def _process_one(doc_ref):
                if self._cancel_requested:
                    return
                async with semaphore:
                    if self._cancel_requested:
                        return
                    max_retries = 3
                    for attempt in range(max_retries):
                        try:
                            async with async_session() as session:
                                stmt = select(OpportunityDocument).where(
                                    OpportunityDocument.id == doc_ref.id
                                )
                                result = await session.execute(stmt)
                                doc = result.scalar_one_or_none()
                                if not doc:
                                    return

                                if doc.download_status == "pending":
                                    await self._download_document(doc, session)

                                if doc.download_status == "downloaded" and doc.ocr_status == "pending":
                                    await self._ocr_document(doc, ocr_settings, session, ocr_client=shared_ocr_client)

                                # Skip classification and embedding for unsupported formats
                                if doc.ocr_status == "skipped":
                                    if doc.classify_status == "pending":
                                        doc.classify_status = "skipped"
                                    if doc.embed_status == "pending":
                                        doc.embed_status = "skipped"

                                # Classify after OCR, before embedding
                                if doc.ocr_status == "completed" and doc.classify_status == "pending":
                                    await self._classify_document(doc, llm_settings, session, llm_client=shared_llm_client)

                                if doc.ocr_status == "completed" and doc.embed_status == "pending":
                                    await self._embed_document(
                                        doc, embed_settings, ocr_settings, session,
                                        embed_client=shared_embed_client,
                                        chroma_collection=shared_chroma_collection,
                                    )

                                # Clear stale error from previous failed attempts
                                all_done = (
                                    doc.download_status == "downloaded"
                                    and doc.ocr_status in ("completed", "skipped")
                                    and doc.classify_status in ("completed", "skipped")
                                    and doc.embed_status in ("completed", "skipped")
                                )
                                if all_done and doc.error_message:
                                    doc.error_message = None

                                await session.commit()
                            break  # Success
                        except Exception as e:
                            err_str = str(e)
                            is_retryable = "1213" in err_str or "Deadlock" in err_str or "1205" in err_str or "Lock wait timeout" in err_str or "PendingRollbackError" in err_str
                            if is_retryable and attempt < max_retries - 1:
                                logger.warning(f"DB contention on document {doc_ref.id}, retry {attempt + 1}: {err_str[:100]}")
                                await asyncio.sleep(0.5 * (attempt + 1))
                                continue
                            logger.error(f"Error processing document {doc_ref.id}: {e}", exc_info=True)
                            self.processing_stats["errors"] += 1
                            try:
                                async with async_session() as err_session:
                                    stmt = select(OpportunityDocument).where(
                                        OpportunityDocument.id == doc_ref.id
                                    )
                                    result = await err_session.execute(stmt)
                                    err_doc = result.scalar_one_or_none()
                                    if err_doc:
                                        err_doc.error_message = str(e)[:2000]
                                        await err_session.commit()
                            except Exception:
                                pass
                            break

                    await self._publish_stats()

            try:
                # Process in batches to allow cancel checks and stats updates
                batch_size = max(num_workers * 3, 10)
                for i in range(0, len(docs), batch_size):
                    if self._cancel_requested:
                        logger.info("Document processing cancelled by user")
                        self.processing_stats["cancelled"] = True
                        break

                    batch = docs[i:i + batch_size]
                    await asyncio.gather(*[_process_one(d) for d in batch])
            finally:
                await shared_ocr_client.aclose()

            self.processing_stats["completed"] = datetime.utcnow().isoformat()
            self.processing_stats["phase"] = "completed"

        except Exception as e:
            logger.error(f"Document processing failed: {e}", exc_info=True)
            self.processing_stats["error"] = str(e)[:500]
        finally:
            self.is_processing = False
            await self._publish_stats()

    async def _download_document(self, doc: OpportunityDocument, session: AsyncSession):
        """Download a single document from Grants.gov."""
        safe_name = re.sub(r'[^\w\-.]', '_', doc.file_name)
        dest_path = os.path.join(
            settings.DOCUMENT_STORAGE_PATH,
            str(doc.opportunity_id),
            f"{doc.attachment_id}_{safe_name}",
        )

        success = await self._grants_client.download_attachment(doc.attachment_id, dest_path)

        if success:
            doc.local_path = dest_path
            doc.download_status = "downloaded"
            self.processing_stats["downloaded"] = self.processing_stats.get("downloaded", 0) + 1
        else:
            doc.download_status = "failed"
            doc.error_message = "Download failed after retries"
            self.processing_stats["errors"] = self.processing_stats.get("errors", 0) + 1

    # File extensions that can be extracted
    _PDF_EXTS = {".pdf"}
    _DOCX_EXTS = {".docx"}
    _DOC_EXTS = {".doc"}
    _XLSX_EXTS = {".xlsx", ".xlsm"}
    _XLS_EXTS = {".xls"}
    _PPTX_EXTS = {".pptx"}
    _HTML_EXTS = {".html", ".htm"}
    _TEXT_EXTS = {".txt", ".csv", ".md"}
    _RTF_EXTS = {".rtf"}
    _UNSUPPORTED_EXTS = {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".exe", ".msi", ".img", ".iso"}

    async def _ocr_document(self, doc: OpportunityDocument, ocr_settings: dict, session: AsyncSession, ocr_client: httpx.AsyncClient | None = None):
        """Extract text from a downloaded document (PDF via OCR, others via direct extraction)."""
        if not doc.local_path or not os.path.exists(doc.local_path):
            doc.ocr_status = "failed"
            doc.error_message = "Local file not found for OCR"
            return

        name = (doc.file_name or "").lower()
        ext = os.path.splitext(name)[1]

        # Route by file extension
        try:
            if ext in self._PDF_EXTS:
                # Try fast local extraction first (pymupdf), fall back to OCR for scanned PDFs
                try:
                    extracted_text = await self._ocr_pymupdf(doc.local_path)
                except Exception as e:
                    logger.warning(f"pymupdf failed for doc {doc.id}: {e}")
                    extracted_text = None
                if extracted_text and len(extracted_text.strip()) >= 200:
                    # pymupdf got enough text — no need for OCR
                    pass
                else:
                    # Scanned/image PDF — send to OCR service (skip files > 5MB to avoid OOM)
                    file_size = doc.file_size or 0
                    if file_size > 5 * 1024 * 1024:
                        logger.info(f"Skipping OCR for large PDF doc {doc.id} ({file_size} bytes), using pymupdf text")
                        if not extracted_text or not extracted_text.strip():
                            doc.ocr_status = "skipped"
                            doc.error_message = "PDF too large for OCR and no extractable text"
                            return
                    else:
                        method = ocr_settings.get("method", "dotsocr")
                        if method == "dotsocr":
                            extracted_text = await self._ocr_dotsocr(doc.local_path, ocr_settings, client=ocr_client)
                        elif method == "pymupdf":
                            pass  # Already tried above
                        else:
                            doc.ocr_status = "failed"
                            doc.error_message = f"Unknown OCR method: {method}"
                            return
            elif ext in self._DOCX_EXTS:
                extracted_text = await self._extract_docx(doc.local_path)
            elif ext in self._DOC_EXTS:
                extracted_text = await self._extract_doc(doc.local_path)
            elif ext in self._XLSX_EXTS:
                extracted_text = await self._extract_xlsx(doc.local_path)
            elif ext in self._XLS_EXTS:
                extracted_text = await self._extract_xls(doc.local_path)
            elif ext in self._PPTX_EXTS:
                extracted_text = await self._extract_pptx(doc.local_path)
            elif ext in self._HTML_EXTS:
                extracted_text = await self._extract_html(doc.local_path)
            elif ext in self._TEXT_EXTS:
                extracted_text = await self._extract_text(doc.local_path)
            elif ext in self._RTF_EXTS:
                extracted_text = await self._extract_rtf(doc.local_path)
            elif ext in self._UNSUPPORTED_EXTS:
                doc.ocr_status = "skipped"
                doc.embed_status = "skipped"
                return
            else:
                # Unknown extension — try plain text, fall back to skip
                try:
                    extracted_text = await self._extract_text(doc.local_path)
                except Exception:
                    doc.ocr_status = "skipped"
                    doc.embed_status = "skipped"
                    return

            if extracted_text:
                doc.ocr_status = "completed"
                doc.extracted_text_length = len(extracted_text)
                text_path = doc.local_path + ".txt"
                with open(text_path, "w", encoding="utf-8") as f:
                    f.write(extracted_text)
                self.processing_stats["ocr_completed"] = self.processing_stats.get("ocr_completed", 0) + 1
            else:
                doc.ocr_status = "failed"
                doc.error_message = "Extraction returned empty text"

        except Exception as e:
            doc.ocr_status = "failed"
            doc.error_message = f"Extraction error: {str(e)[:500]}"
            logger.error(f"Text extraction failed for document {doc.id}: {e}", exc_info=True)

    async def _ocr_dotsocr(self, file_path: str, ocr_settings: dict, client: httpx.AsyncClient | None = None) -> str | None:
        """Send PDF to dotsocr endpoint, receive markdown."""
        endpoint = ocr_settings.get("endpoint_url", settings.OCR_ENDPOINT_URL)

        should_close = client is None
        if client is None:
            client = httpx.AsyncClient(timeout=300.0)
        try:
            with open(file_path, "rb") as f:
                files = {"file": (os.path.basename(file_path), f, "application/pdf")}
                response = await client.post(endpoint, files=files)

            response.raise_for_status()

            # dotsocr returns markdown text
            content_type = response.headers.get("content-type", "")
            if "json" in content_type:
                data = response.json()
                return data.get("text") or data.get("markdown") or data.get("content", "")
            else:
                return response.text
        finally:
            if should_close:
                await client.aclose()

    async def _ocr_pymupdf(self, file_path: str) -> str | None:
        """Extract text from PDF using PyMuPDF (local, no network)."""
        loop = asyncio.get_event_loop()

        def _extract():
            import pymupdf
            doc = pymupdf.open(file_path)
            pages = []
            for page in doc:
                pages.append(page.get_text())
            doc.close()
            return "\n\n".join(pages)

        return await loop.run_in_executor(None, _extract)

    async def _extract_docx(self, file_path: str) -> str | None:
        """Extract text from .docx using python-docx."""
        loop = asyncio.get_event_loop()

        def _extract():
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)

        return await loop.run_in_executor(None, _extract)

    async def _extract_doc(self, file_path: str) -> str | None:
        """Extract text from legacy .doc using antiword."""
        proc = await asyncio.create_subprocess_exec(
            "antiword", file_path,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await proc.communicate()
        if proc.returncode != 0:
            raise RuntimeError(f"antiword failed: {stderr.decode(errors='replace')[:200]}")
        return stdout.decode("utf-8", errors="replace")

    async def _extract_xlsx(self, file_path: str) -> str | None:
        """Extract text from .xlsx/.xlsm using openpyxl."""
        loop = asyncio.get_event_loop()

        def _extract():
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append("\t".join(cells))
                if rows:
                    parts.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
            wb.close()
            return "\n\n".join(parts)

        return await loop.run_in_executor(None, _extract)

    async def _extract_xls(self, file_path: str) -> str | None:
        """Extract text from legacy .xls using xlrd."""
        loop = asyncio.get_event_loop()

        def _extract():
            import xlrd
            wb = xlrd.open_workbook(file_path)
            parts = []
            for sheet in wb.sheets():
                rows = []
                for rx in range(sheet.nrows):
                    cells = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
                    if any(cells):
                        rows.append("\t".join(cells))
                if rows:
                    parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
            return "\n\n".join(parts)

        return await loop.run_in_executor(None, _extract)

    async def _extract_pptx(self, file_path: str) -> str | None:
        """Extract text from .pptx using python-pptx."""
        loop = asyncio.get_event_loop()

        def _extract():
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                texts.append(para.text.strip())
                if texts:
                    parts.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
            return "\n\n".join(parts)

        return await loop.run_in_executor(None, _extract)

    async def _extract_html(self, file_path: str) -> str | None:
        """Extract text from .html/.htm using BeautifulSoup."""
        loop = asyncio.get_event_loop()

        def _extract():
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            return soup.get_text(separator="\n", strip=True)

        return await loop.run_in_executor(None, _extract)

    async def _extract_text(self, file_path: str) -> str | None:
        """Read plain text files (.txt, .csv, .md)."""
        loop = asyncio.get_event_loop()

        def _extract():
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        return await loop.run_in_executor(None, _extract)

    async def _extract_rtf(self, file_path: str) -> str | None:
        """Extract text from .rtf using striprtf."""
        loop = asyncio.get_event_loop()

        def _extract():
            from striprtf.striprtf import rtf_to_text
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return rtf_to_text(f.read())

        return await loop.run_in_executor(None, _extract)

    _CLASSIFICATION_CATEGORIES = [
        "solicitation",      # The primary solicitation — RFP, RFA, RFQ, BAA, NOFO, FOA
        "budget_template",   # Budget forms, templates, examples
        "application_form",  # SF424, fillable application forms
        "legal_compliance",  # Privacy Act statements, certifications, representations
        "amendment",         # Amendments, modifications to original solicitation
        "instructions",      # How-to guides, submission instructions, README
        "presentation",      # Slide decks, briefings
        "data_spreadsheet",  # Data files, spreadsheets, tracking sheets
        "other",             # Anything that doesn't fit above
    ]

    # --- Heuristic pre-classification rules ---
    _HEURISTIC_RULES = {
        "solicitation": {
            "filename_patterns": [
                r'\brfp\b', r'\brfa\b', r'\brfq\b', r'\bbaa\b',
                r'\bsolicitation\b', r'\bfoa\b', r'\bnofo\b',
                r'notice.*funding.*opportunity', r'funding.*opportunity.*announcement',
                r'\bprogram\s*announcement\b', r'\bbroad\s*agency\s*announcement\b',
            ],
            "folder_patterns": [
                r'full\s*announcement', r'solicitation', r'\bbaa\b',
                r'\brfa\b', r'\brfp\b', r'\bnofo\b', r'notice.*funding',
                r'funding.*opportunity', r'\bfoa\b',
            ],
        },
        "amendment": {
            "filename_patterns": [r'\bamendment\b', r'\bmodification\b', r'\baddendum\b'],
            "folder_patterns": [r'\bamendment\b', r'\bmodification\b'],
        },
        "budget_template": {
            "filename_patterns": [r'\bbudget\b.*\btemplate\b', r'\bbudget\b.*\bform\b', r'\bsf[\-_]?424a\b'],
            "folder_patterns": [r'\bbudget\b'],
        },
        "application_form": {
            "filename_patterns": [r'\bsf[\-_]?424\b', r'\bapplication\b.*\bform\b', r'\brr_.*form\b'],
            "folder_patterns": [r'\bapplication\b.*\bform\b', r'\brequired\s*form\b'],
        },
        "instructions": {
            "filename_patterns": [r'\binstructions?\b', r'\bapplication\b.*\bguide\b', r'\breadme\b'],
            "folder_patterns": [r'\binstructions?\b', r'\bguide\b'],
        },
    }

    @classmethod
    def _heuristic_classify(cls, file_name: str, folder_name: str | None) -> str | None:
        """Try to classify a document based on filename/folder patterns.

        Returns category string if a confident match is found, None otherwise.
        """
        fn_lower = (file_name or "").lower()
        folder_lower = (folder_name or "").lower()

        for category, rules in cls._HEURISTIC_RULES.items():
            # Check filename patterns
            for pattern in rules.get("filename_patterns", []):
                if re.search(pattern, fn_lower):
                    return category
            # Check folder patterns
            for pattern in rules.get("folder_patterns", []):
                if folder_lower and re.search(pattern, folder_lower):
                    return category

        return None

    _CLASSIFY_SYSTEM_PROMPT = """You are a document classifier for federal grant opportunities. Given the beginning of a document, classify it into exactly one category.

Categories:
- solicitation: The primary solicitation document — RFP, RFA, RFQ, BAA, NOFO, FOA, or similar. This is the main document describing what is being solicited and how to apply.
- budget_template: Budget forms, budget templates, cost examples
- application_form: SF424, fillable application forms, registration forms
- legal_compliance: Privacy Act statements, certifications, representations, compliance documents
- amendment: Amendments or modifications to the original solicitation
- instructions: How-to guides, submission instructions, README files, application guides
- presentation: Slide decks, briefings, webinar materials
- data_spreadsheet: Data files, spreadsheets, tracking sheets
- other: Anything that doesn't fit the above categories

Examples:
- "FY2025_BAA_DARPA-PA-25-01.pdf" in folder "Full Announcement" → solicitation
- "NOFO-HHS-2025-001.pdf" in folder "NOFO" → solicitation
- "RFA-CA-25-003.pdf" in folder "Full Announcement" → solicitation
- "SF424_RR_Budget.pdf" in folder "Required Forms" → application_form
- "Budget_Justification_Template.docx" in folder "Budget" → budget_template
- "Amendment_003.pdf" in folder "Amendments" → amendment
- "Application_Guide.pdf" in folder "Instructions" → instructions

Respond with ONLY a JSON object: {"category": "<category>", "confidence": "high|medium|low"}"""

    async def _classify_document(
        self,
        doc: OpportunityDocument,
        llm_settings: dict,
        session: AsyncSession,
        llm_client=None,
    ):
        """Classify a document using heuristics first, then LLM fallback."""
        # --- Phase 1: Try heuristic classification ---
        heuristic_result = self._heuristic_classify(doc.file_name, doc.folder_name)
        if heuristic_result:
            doc.doc_category = heuristic_result
            doc.classify_status = "completed"
            self.processing_stats["classified"] = self.processing_stats.get("classified", 0) + 1
            return

        # --- Phase 2: LLM-based classification ---
        text_path = doc.local_path + ".txt"
        if not os.path.exists(text_path):
            doc.classify_status = "failed"
            doc.error_message = "Extracted text file not found for classification"
            return

        # Read first ~3000 tokens worth of text (rough estimate: 4 chars per token)
        try:
            with open(text_path, "r", encoding="utf-8") as f:
                preview_text = f.read(12000)
        except Exception as e:
            doc.classify_status = "failed"
            doc.error_message = f"Classification read error: {str(e)[:200]}"
            return

        if not preview_text.strip():
            doc.doc_category = "other"
            doc.classify_status = "completed"
            return

        base_url = llm_settings.get("base_url", "")
        model = llm_settings.get("model", "")
        api_key = llm_settings.get("api_key", "")

        if not base_url or not model:
            # No LLM configured — skip classification silently
            doc.doc_category = "other"
            doc.classify_status = "skipped"
            return

        try:
            if llm_client is None:
                from openai import AsyncOpenAI
                llm_client = AsyncOpenAI(base_url=base_url, api_key=api_key or "not-needed")

            folder_line = f"\nFolder: {doc.folder_name}" if doc.folder_name else ""
            user_msg = f"File name: {doc.file_name}{folder_line}\n\nDocument text (first ~3000 tokens):\n{preview_text}"

            response = await llm_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": self._CLASSIFY_SYSTEM_PROMPT},
                    {"role": "user", "content": user_msg},
                ],
                max_tokens=60,
                temperature=0.0,
                timeout=30,
            )

            raw = (response.choices[0].message.content or "").strip()

            # Parse JSON — strip markdown fences if present
            import json
            cleaned = raw
            if "```" in cleaned:
                cleaned = re.sub(r"```(?:json)?\s*", "", cleaned)
                cleaned = cleaned.replace("```", "").strip()
            # Try to extract JSON object
            match = re.search(r'\{[^}]+\}', cleaned)
            if match:
                parsed = json.loads(match.group())
            else:
                parsed = json.loads(cleaned)

            category = parsed.get("category", "other").strip().lower()
            # Accept legacy category names from LLM and map them
            if category in ("rfp_rfa", "nofo"):
                category = "solicitation"
            if category not in self._CLASSIFICATION_CATEGORIES:
                category = "other"

            doc.doc_category = category
            doc.classify_status = "completed"
            self.processing_stats["classified"] = self.processing_stats.get("classified", 0) + 1

        except Exception as e:
            logger.warning(f"Classification failed for document {doc.id}, defaulting to 'other': {e}")
            doc.doc_category = "other"
            doc.classify_status = "completed"
            self.processing_stats["classified"] = self.processing_stats.get("classified", 0) + 1

    async def _embed_document(self, doc: OpportunityDocument, embed_settings: dict, ocr_settings: dict, session: AsyncSession, embed_client=None, chroma_collection=None):
        """Chunk extracted text, generate embeddings, store in ChromaDB."""
        text_path = doc.local_path + ".txt"
        if not os.path.exists(text_path):
            doc.embed_status = "failed"
            doc.error_message = "Extracted text file not found"
            return

        with open(text_path, "r", encoding="utf-8") as f:
            full_text = f.read()

        if not full_text.strip():
            doc.embed_status = "failed"
            doc.error_message = "Extracted text is empty"
            return

        base_url = embed_settings.get("base_url", "")
        model = embed_settings.get("model", "")
        api_key = embed_settings.get("api_key", "")

        if not base_url or not model:
            doc.embed_status = "failed"
            doc.error_message = "Embedding endpoint not configured"
            return

        try:
            # Chunk the text using token-based sizing
            chunk_size = ocr_settings.get("chunk_size_tokens", 1000)
            chunk_overlap = ocr_settings.get("chunk_overlap_tokens", 200)
            chunks = self._chunk_text(full_text, chunk_size=chunk_size, overlap=chunk_overlap)

            # Delete existing chunks for this document
            await session.execute(
                text("DELETE FROM document_chunks WHERE document_id = :doc_id"),
                {"doc_id": doc.id},
            )

            # Generate embeddings in batches — use shared client if provided
            if embed_client is None:
                from openai import AsyncOpenAI
                embed_client = AsyncOpenAI(base_url=base_url, api_key=api_key or "not-needed")

            all_embeddings = []
            batch_size = 20
            for i in range(0, len(chunks), batch_size):
                batch_texts = [c["text"] for c in chunks[i:i + batch_size]]
                response = await embed_client.embeddings.create(
                    model=model,
                    input=batch_texts,
                    timeout=60,
                )
                all_embeddings.extend([d.embedding for d in response.data])

            # Store in ChromaDB — use shared collection if provided
            if chroma_collection is None:
                import chromadb
                chroma_client = chromadb.HttpClient(
                    host=settings.CHROMADB_HOST,
                    port=settings.CHROMADB_PORT,
                )
                chroma_collection = chroma_client.get_or_create_collection(
                    name="opportunity_documents",
                    metadata={"hnsw:space": "cosine"},
                )
            collection = chroma_collection

            chroma_ids = []
            chroma_embeddings = []
            chroma_documents = []
            chroma_metadatas = []

            # Build opportunity metadata for ChromaDB enrichment
            opp_meta = {}
            if doc.opportunity:
                opp_meta = {
                    "opportunity_number": doc.opportunity.opportunity_number or "",
                    "opportunity_title": doc.opportunity.title or "",
                    "agency_code": doc.opportunity.agency_code or "",
                    "grants_gov_url": doc.opportunity.grants_gov_url or "",
                }

            chunk_rows = []
            for idx, chunk_data in enumerate(chunks):
                chroma_id = f"doc_{doc.id}_chunk_{idx}_{uuid.uuid4().hex[:8]}"
                chroma_ids.append(chroma_id)
                chroma_embeddings.append(all_embeddings[idx])
                chroma_documents.append(chunk_data["text"])
                chroma_metadatas.append({
                    "document_id": doc.id,
                    "opportunity_id": doc.opportunity_id,
                    "attachment_id": doc.attachment_id,
                    "chunk_index": idx,
                    "total_chunks": len(chunks),
                    "file_name": doc.file_name,
                    "doc_category": doc.doc_category or "",
                    "source": doc.source or "",
                    "local_path": doc.local_path or "",
                    **opp_meta,
                })

                chunk_row = DocumentChunk(
                    document_id=doc.id,
                    chunk_index=idx,
                    chunk_text=chunk_data["text"],
                    char_offset=chunk_data["offset"],
                    char_length=chunk_data["length"],
                    chroma_id=chroma_id,
                )
                chunk_rows.append(chunk_row)

            # Upsert to ChromaDB in batches
            for i in range(0, len(chroma_ids), 100):
                collection.upsert(
                    ids=chroma_ids[i:i + 100],
                    embeddings=chroma_embeddings[i:i + 100],
                    documents=chroma_documents[i:i + 100],
                    metadatas=chroma_metadatas[i:i + 100],
                )

            # Save chunk rows to MariaDB
            for row in chunk_rows:
                session.add(row)
            await session.flush()

            doc.embed_status = "completed"
            doc.chunk_count = len(chunks)
            self.processing_stats["embedded"] = self.processing_stats.get("embedded", 0) + 1

        except Exception as e:
            doc.embed_status = "failed"
            doc.error_message = f"Embedding error: {str(e)[:500]}"
            logger.error(f"Embedding failed for document {doc.id}: {e}", exc_info=True)

    def _chunk_text(
        self, text_content: str, chunk_size: int = 1000, overlap: int = 200
    ) -> list[dict]:
        """Split text at paragraph/sentence boundaries using token counts.

        Args:
            text_content: Full text to chunk.
            chunk_size: Target chunk size in tokens.
            overlap: Overlap between chunks in tokens.

        Returns list of {text, offset, length}.
        """
        if not text_content:
            return []

        import tiktoken
        try:
            enc = tiktoken.get_encoding("cl100k_base")
        except Exception:
            enc = tiktoken.get_encoding("gpt2")

        def _token_len(s: str) -> int:
            return len(enc.encode(s, disallowed_special=()))

        # Split into paragraphs
        paragraphs = re.split(r'\n\s*\n', text_content)

        chunks = []
        current_paras: list[str] = []
        current_tokens = 0
        current_offset = 0
        pos = 0

        def _split_oversized(paragraph: str, max_tokens: int) -> list[str]:
            """Split a paragraph that exceeds max_tokens into smaller pieces."""
            # Try sentence splitting first
            sentences = re.split(r'(?<=[.!?])\s+', paragraph)
            if len(sentences) <= 1:
                # Single sentence — hard-truncate by tokens
                tokens = enc.encode(paragraph, disallowed_special=())
                result = []
                for i in range(0, len(tokens), max_tokens):
                    result.append(enc.decode(tokens[i:i + max_tokens]))
                return result

            pieces = []
            current = []
            current_tok = 0
            for sent in sentences:
                sent_tok = _token_len(sent)
                if sent_tok > max_tokens:
                    # Flush current
                    if current:
                        pieces.append(" ".join(current))
                        current = []
                        current_tok = 0
                    # Hard-truncate the oversized sentence
                    tokens = enc.encode(sent, disallowed_special=())
                    for i in range(0, len(tokens), max_tokens):
                        pieces.append(enc.decode(tokens[i:i + max_tokens]))
                elif current_tok + sent_tok > max_tokens and current:
                    pieces.append(" ".join(current))
                    current = [sent]
                    current_tok = sent_tok
                else:
                    current.append(sent)
                    current_tok += sent_tok
            if current:
                pieces.append(" ".join(current))
            return pieces

        for para in paragraphs:
            para = para.strip()
            if not para:
                pos += 2
                continue

            para_start = text_content.find(para, pos)
            if para_start == -1:
                para_start = pos
            pos = para_start + len(para)

            para_tokens = _token_len(para)

            # If a single paragraph exceeds chunk_size, split it
            if para_tokens > chunk_size:
                # Flush current chunk first
                if current_paras:
                    chunk_text = "\n\n".join(current_paras).strip()
                    chunks.append({
                        "text": chunk_text,
                        "offset": current_offset,
                        "length": len(chunk_text),
                    })
                    current_paras = []
                    current_tokens = 0

                sub_parts = _split_oversized(para, chunk_size)
                sub_offset = para_start
                for sp in sub_parts:
                    chunks.append({
                        "text": sp,
                        "offset": sub_offset,
                        "length": len(sp),
                    })
                    sub_offset += len(sp)
                current_offset = pos
                continue

            if current_tokens + para_tokens > chunk_size and current_paras:
                chunk_text = "\n\n".join(current_paras).strip()
                chunks.append({
                    "text": chunk_text,
                    "offset": current_offset,
                    "length": len(chunk_text),
                })

                # Build overlap from trailing paragraphs
                if overlap > 0:
                    overlap_paras: list[str] = []
                    overlap_tokens = 0
                    for p in reversed(current_paras):
                        p_tok = _token_len(p)
                        if overlap_tokens + p_tok > overlap and overlap_paras:
                            break
                        overlap_paras.insert(0, p)
                        overlap_tokens += p_tok
                    current_paras = overlap_paras + [para]
                    current_tokens = overlap_tokens + para_tokens
                    # Approximate offset from overlap start
                    overlap_text = "\n\n".join(overlap_paras)
                    current_offset = para_start - len(overlap_text) - 2 if overlap_text else para_start
                else:
                    current_paras = [para]
                    current_tokens = para_tokens
                    current_offset = para_start
            else:
                if not current_paras:
                    current_offset = para_start
                current_paras.append(para)
                current_tokens += para_tokens

        # Final chunk
        if current_paras:
            chunk_text = "\n\n".join(current_paras).strip()
            if chunk_text:
                chunks.append({
                    "text": chunk_text,
                    "offset": current_offset,
                    "length": len(chunk_text),
                })

        return chunks

    async def semantic_search(
        self,
        query: str,
        n_results: int = 10,
        opportunity_id: int | None = None,
    ) -> list[dict]:
        """Search ChromaDB for relevant document chunks.

        Returns list of {chunk_text, document_id, opportunity_id, file_name, score}.
        """
        async with async_session() as session:
            embed_settings = await settings_service.get_embedding_settings(session)

        base_url = embed_settings.get("base_url", "")
        model = embed_settings.get("model", "")
        api_key = embed_settings.get("api_key", "")

        if not base_url or not model:
            return []

        try:
            from openai import AsyncOpenAI
            embed_client = AsyncOpenAI(base_url=base_url, api_key=api_key or "not-needed")
            response = await embed_client.embeddings.create(
                model=model,
                input=query,
                timeout=30,
            )
            query_embedding = response.data[0].embedding

            import chromadb
            chroma_client = chromadb.HttpClient(
                host=settings.CHROMADB_HOST,
                port=settings.CHROMADB_PORT,
            )
            collection = chroma_client.get_or_create_collection(
                name="opportunity_documents",
                metadata={"hnsw:space": "cosine"},
            )

            where_filter = None
            if opportunity_id is not None:
                where_filter = {"opportunity_id": opportunity_id}

            results = collection.query(
                query_embeddings=[query_embedding],
                n_results=n_results,
                where=where_filter,
                include=["documents", "metadatas", "distances"],
            )

            items = []
            if results and results.get("documents"):
                docs = results["documents"][0]
                metas = results["metadatas"][0]
                distances = results["distances"][0]

                for doc_text, meta, dist in zip(docs, metas, distances):
                    items.append({
                        "chunk_text": doc_text,
                        "document_id": meta.get("document_id"),
                        "opportunity_id": meta.get("opportunity_id"),
                        "file_name": meta.get("file_name"),
                        "score": 1.0 - dist,  # cosine distance to similarity
                    })

            return items

        except Exception as e:
            logger.error(f"Semantic search failed: {e}", exc_info=True)
            return []

    def cancel_processing(self):
        """Request cancellation of document processing."""
        self._cancel_requested = True

    async def get_processing_status(self) -> dict:
        """Get current processing status. Redis is the single source of truth."""
        # Check the simple atomic flag first — most reliable
        is_active = False
        flag = None
        try:
            flag = await cache_service._redis.get(REDIS_DOC_PROCESSING_FLAG)
            is_active = flag == b"1" or flag == "1"
        except Exception as e:
            logger.warning(f"Redis flag check failed: {e}, falling back to local state")
            is_active = self.is_processing

        # Get detailed stats from Redis (or local fallback)
        shared = await self._get_shared_stats()
        stats = shared if shared else self.processing_stats

        # If idle, ensure we have the completion timestamp from persistent key
        if not is_active and not stats.get("completed"):
            try:
                completed = await cache_service._redis.get(REDIS_DOC_COMPLETED_KEY)
                if completed:
                    stats = {**stats, "completed": completed}
            except Exception:
                pass

        return {"is_processing": is_active, "stats": stats}

    async def get_document_counts(self) -> dict:
        """Get aggregate counts for the admin dashboard (open opportunities only)."""
        async with async_session() as session:
            today = date.today()
            base = (
                select(func.count(OpportunityDocument.id))
                .join(Opportunity, OpportunityDocument.opportunity_id == Opportunity.id)
                .where(
                    Opportunity.status != "archived",
                    or_(Opportunity.close_date >= today, Opportunity.close_date.is_(None)),
                )
            )
            total = await session.execute(base)
            downloaded = await session.execute(
                base.where(OpportunityDocument.download_status == "downloaded")
            )
            ocr_completed = await session.execute(
                base.where(OpportunityDocument.ocr_status == "completed")
            )
            classified = await session.execute(
                base.where(OpportunityDocument.classify_status == "completed")
            )
            embedded = await session.execute(
                base.where(OpportunityDocument.embed_status == "completed")
            )
            errors = await session.execute(
                base.where(or_(
                    OpportunityDocument.download_status == "failed",
                    OpportunityDocument.ocr_status == "failed",
                    OpportunityDocument.embed_status == "failed",
                ))
            )
            pending = await session.execute(
                base.where(or_(
                    OpportunityDocument.download_status == "pending",
                    OpportunityDocument.ocr_status == "pending",
                    OpportunityDocument.classify_status == "pending",
                    OpportunityDocument.embed_status == "pending",
                ))
            )

            return {
                "total": total.scalar() or 0,
                "downloaded": downloaded.scalar() or 0,
                "ocr_completed": ocr_completed.scalar() or 0,
                "classified": classified.scalar() or 0,
                "embedded": embedded.scalar() or 0,
                "errors": errors.scalar() or 0,
                "pending": pending.scalar() or 0,
            }

    async def get_recent_errors(self, limit: int = 20) -> list[dict]:
        """Get recent document processing errors for the admin UI."""
        async with async_session() as session:
            stmt = (
                select(OpportunityDocument)
                .where(OpportunityDocument.error_message.is_not(None))
                .order_by(OpportunityDocument.updated_at.desc())
                .limit(limit)
            )
            result = await session.execute(stmt)
            docs = result.scalars().all()

            return [
                {
                    "id": doc.id,
                    "file_name": doc.file_name,
                    "attachment_id": doc.attachment_id,
                    "download_status": doc.download_status,
                    "ocr_status": doc.ocr_status,
                    "embed_status": doc.embed_status,
                    "error_message": doc.error_message,
                    "updated_at": doc.updated_at.isoformat() if doc.updated_at else "",
                }
                for doc in docs
            ]

    async def _publish_stats(self):
        """Publish processing stats to Redis for cross-worker visibility."""
        try:
            import json
            payload = {**self.processing_stats, "is_processing": self.is_processing}
            pipe = cache_service._redis.pipeline()
            pipe.set(REDIS_DOC_SYNC_KEY, json.dumps(payload), ex=300)
            # Separate atomic flag — simple string, no JSON
            if self.is_processing:
                pipe.set(REDIS_DOC_PROCESSING_FLAG, "1", ex=300)
            else:
                pipe.delete(REDIS_DOC_PROCESSING_FLAG)
                # Store completion time persistently (24h TTL)
                if payload.get("completed"):
                    pipe.set(REDIS_DOC_COMPLETED_KEY, payload["completed"], ex=86400)
            await pipe.execute()
        except Exception:
            pass

    async def _get_shared_stats(self) -> dict | None:
        """Read processing stats from Redis."""
        try:
            import json
            data = await cache_service._redis.get(REDIS_DOC_SYNC_KEY)
            if data:
                return json.loads(data)
        except Exception:
            pass
        return None


document_service = DocumentService()
